import os
import io
import tempfile
from groq import Groq

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import fitz
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def _get_client():
    api_key = os.environ.get("GROQ_API_KEY", "")
    if not api_key:
        raise RuntimeError("Clé GROQ_API_KEY manquante dans .streamlit/secrets.toml")
    return Groq(api_key=api_key)


def _translate_with_groq(text: str, source_lang: str, target_lang: str) -> str:
    if not text.strip():
        return text
    client   = _get_client()
    prompt   = f"Translate from {source_lang} to {target_lang}. Return ONLY the translated text, no explanation:\n\n{text}"
    response = client.chat.completions.create(
        model       = "llama-3.3-70b-versatile",
        messages    = [{"role": "user", "content": prompt}],
        max_tokens  = 2048,
        temperature = 0.1
    )
    return response.choices[0].message.content.strip()


def _translate_chunks(text: str, source_lang: str, target_lang: str,
                       chunk_size: int = 2000) -> str:
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    return "\n".join(_translate_with_groq(c, source_lang, target_lang)
                     for c in chunks if c.strip())


class DocumentTranslator:
    SUPPORTED_EXTENSIONS = {".txt", ".docx", ".pdf", ".pptx"}

    def translate(self, file, target_lang: str = "en", source_lang: str = "fr"):
        filename  = getattr(file, "name", None) or getattr(file, "filename", "fichier.txt")
        extension = os.path.splitext(filename.lower())[1]
        if extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Type non supporté : {extension}")
        file_bytes = file.read() if hasattr(file, "read") else b""
        if extension == ".docx":
            return self._translate_docx(file_bytes, source_lang, target_lang)
        elif extension == ".txt":
            return self._translate_txt(file_bytes, source_lang, target_lang)
        elif extension == ".pdf":
            return self._translate_pdf(file_bytes, source_lang, target_lang)
        elif extension == ".pptx":
            return self._translate_pptx(file_bytes, source_lang, target_lang)

    def _translate_docx(self, file_bytes, source_lang, target_lang):
        if not DOCX_AVAILABLE:
            raise ImportError("Installez : pip install python-docx")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
        finally:
            os.unlink(tmp_path)
        for paragraph in doc.paragraphs:
            if not paragraph.text.strip():
                continue
            translated = _translate_with_groq(paragraph.text, source_lang, target_lang)
            if paragraph.runs:
                paragraph.runs[0].text = translated
                for run in paragraph.runs[1:]:
                    run.text = ""
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():
                            translated = _translate_with_groq(
                                paragraph.text, source_lang, target_lang)
                            if paragraph.runs:
                                paragraph.runs[0].text = translated
                                for run in paragraph.runs[1:]:
                                    run.text = ""
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        return (output.getvalue(), ".docx",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    def _translate_pptx(self, file_bytes, source_lang, target_lang):
        if not PPTX_AVAILABLE:
            raise ImportError("Installez : pip install python-pptx")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
        finally:
            os.unlink(tmp_path)
        # Traduit chaque slide en conservant la mise en forme
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            run.text = _translate_with_groq(
                                run.text, source_lang, target_lang)
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return (output.getvalue(), ".pptx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation")

    def _translate_txt(self, file_bytes, source_lang, target_lang):
        text       = file_bytes.decode("utf-8", errors="ignore")
        translated = _translate_chunks(text, source_lang, target_lang)
        return translated.encode("utf-8"), ".txt", "text/plain"

    def _translate_pdf(self, file_bytes, source_lang, target_lang):
        if not PDF_AVAILABLE:
            raise ImportError("Installez : pip install pymupdf")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            text = ""
            with fitz.open(tmp_path) as doc:
                for page in doc:
                    text += page.get_text()
        finally:
            os.unlink(tmp_path)
        if not text.strip():
            raise ValueError("Impossible d'extraire du texte du PDF.")
        translated = _translate_chunks(text, source_lang, target_lang)
        return translated.encode("utf-8"), ".txt", "text/plain"