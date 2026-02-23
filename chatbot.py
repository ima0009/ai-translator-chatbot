import os
import time
import base64
from groq import Groq
from language_detector import LanguageDetector

CHAT_MODEL   = "llama-3.3-70b-versatile"
VISION_MODEL = "meta-llama/llama-4-scout-17b-16e-instruct"

SYSTEM_PROMPT = """Tu es MERCREDI, un assistant pédagogique intelligent et multidomaines, conçu pour aider les étudiants de tous niveaux.

Tes domaines d'expertise :

🌍 LANGUES & TRADUCTION
- Traduction précise et naturelle avec nuances et alternatives
- Grammaire, conjugaison, orthographe avec explications claires
- Vocabulaire avancé, expressions idiomatiques, faux amis
- Préparation aux examens : DELF, DALF, TOEFL, IELTS, Cambridge

📚 MATIÈRES ACADÉMIQUES
- Mathématiques : algèbre, analyse, géométrie, statistiques
- Sciences : physique, chimie, biologie, SVT
- Histoire & Géographie : événements, analyses, cartographie
- Philosophie : courants, auteurs, dissertation
- Économie & Gestion : concepts, exercices, cas pratiques
- Informatique : programmation, algorithmes, bases de données
- Droit : principes juridiques, cas pratiques

✍️ RÉDACTION & MÉTHODOLOGIE
- Rédaction de dissertations, rapports, mémoires
- Résumés, synthèses, analyses de textes
- Correction et amélioration de textes
- Citations, bibliographies, normes APA, MLA

🖼️ ANALYSE VISUELLE & DOCUMENTS
- Analyse d'images, schémas, graphiques, tableaux
- Lecture et analyse de documents (.txt, .pdf, .docx, .pptx)
- Résumé, correction, explication de contenu

🧠 MÉTHODES DE TRAVAIL
- Techniques de mémorisation et révision
- Gestion du temps et organisation
- Conseils pour les examens et concours

Règles :
1. Réponds TOUJOURS dans la même langue que l'étudiant
2. Sois pédagogue : explique le POURQUOI avec des exemples concrets
3. Structure tes réponses clairement pour les sujets complexes
4. Encourage l'étudiant et valorise ses efforts
5. Si tu n'es pas sûr d'une réponse, dis-le honnêtement"""


def _get_client():
    api_key = os.environ.get("GROQ_API_KEY", "")
    if not api_key:
        raise RuntimeError("❌ Clé GROQ_API_KEY manquante dans .streamlit/secrets.toml")
    return Groq(api_key=api_key)


def _extract_doc_text(file_bytes: bytes, filename: str) -> str:
    """Extrait le texte d'un document."""
    ext = os.path.splitext(filename.lower())[1]

    if ext == ".txt":
        return file_bytes.decode("utf-8", errors="ignore")

    elif ext == ".pdf":
        import fitz
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            text = ""
            with fitz.open(tmp_path) as doc:
                for page in doc:
                    text += page.get_text()
        finally:
            os.unlink(tmp_path)
        return text

    elif ext == ".docx":
        from docx import Document
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc  = Document(tmp_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        finally:
            os.unlink(tmp_path)
        return text

    elif ext == ".pptx":
        from pptx import Presentation
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs  = Presentation(tmp_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text_frame.text + "\n"
        finally:
            os.unlink(tmp_path)
        return text

    return ""


def _image_to_base64(image_bytes: bytes, mime_type: str = "image/jpeg") -> str:
    return f"data:{mime_type};base64,{base64.b64encode(image_bytes).decode('utf-8')}"


class Chatbot:
    def __init__(self):
        self.detector             = LanguageDetector()
        self.conversation_history = []

    def respond(self, message: str,
                image_bytes: bytes = None, image_mime: str = "image/jpeg",
                doc_bytes: bytes = None, doc_name: str = None,
                on_retry=None) -> str:

        # Cas 1 : document texte joint
        if doc_bytes and doc_name:
            doc_text = _extract_doc_text(doc_bytes, doc_name)
            if doc_text.strip():
                # Tronquer si trop long
                doc_text  = doc_text[:6000]
                user_content = f"[Document : {doc_name}]\n\n{doc_text}\n\n---\n{message if message.strip() else 'Analyse ce document.'}"
            else:
                user_content = message
            model = CHAT_MODEL
            # Ajouter à l'historique avec un string simple
            self.conversation_history.append({"role": "user", "content": user_content})

        # Cas 2 : image jointe
        elif image_bytes:
            image_url    = _image_to_base64(image_bytes, image_mime)
            # Pour l'API, on utilise le format multimodal
            user_content_multimodal = [
                {"type": "image_url", "image_url": {"url": image_url}},
                {"type": "text",      "text": message if message.strip() 
                                          else "Analyse cette image et explique son contenu."}
            ]
            model = VISION_MODEL
            # Pour l'historique, on stocke une version texte simplifiée
            text_version = message if message.strip() else "[Image envoyée]"
            self.conversation_history.append({"role": "user", "content": text_version})
            # On utilise le contenu multimodal pour l'appel API, pas pour l'historique
            user_content_for_api = user_content_multimodal

        # Cas 3 : texte seul
        else:
            user_content = message
            model        = CHAT_MODEL
            self.conversation_history.append({"role": "user", "content": user_content})
            user_content_for_api = user_content

        try:
            client = _get_client()

            # Construction des messages pour l'API
            if image_bytes:
                # Pour les images : pas d'historique, juste le message multimodal
                messages = [
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user",   "content": user_content_for_api}
                ]
            else:
                # Pour le texte seul : on utilise l'historique complet
                messages = [
                    {"role": "system", "content": SYSTEM_PROMPT},
                    *self.conversation_history
                ]

            retries = 3
            wait    = 3
            for attempt in range(retries):
                try:
                    response = client.chat.completions.create(
                        model       = model,
                        messages    = messages,
                        temperature = 0.3,
                        max_tokens  = 2048,
                    )
                    break
                except Exception as e:
                    if "429" in str(e) and attempt < retries - 1:
                        if on_retry:
                            on_retry(attempt + 1, wait)
                        time.sleep(wait)
                    else:
                        raise e

            reply = response.choices[0].message.content

        except RuntimeError as e:
            return str(e)
        except Exception as e:
            return f"⚠️ Erreur : {str(e)}"

        self.conversation_history.append({"role": "assistant", "content": reply})
        if len(self.conversation_history) > 40:
            self.conversation_history = self.conversation_history[-40:]

        return reply

    def reset_conversation(self):
        self.conversation_history = []
        return "Conversation réinitialisée."